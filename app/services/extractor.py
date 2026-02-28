"""
Document extraction service.
Extracts text from PDF, DOCX, and PPTX files.
Returns a list of page strings.

Libraries:
  PDF  → PyMuPDF (fitz)
  DOCX → python-docx
  PPTX → python-pptx
"""
import io
from typing import List


def extract_pdf(data: bytes) -> List[str]:
    """Extract text from PDF bytes. Returns one string per page."""
    try:
        import fitz  # PyMuPDF

        doc  = fitz.open(stream=data, filetype="pdf")
        pages = []
        for page in doc:
            text = page.get_text("text").strip()
            if text:
                pages.append(text)
            else:
                pages.append("[This page contains an image or diagram.]")
        doc.close()
        return pages if pages else ["[No text content found in this PDF.]"]
    except Exception as e:
        return [f"[Error extracting PDF: {str(e)}]"]


def extract_docx(data: bytes) -> List[str]:
    """Extract text from DOCX bytes. Groups paragraphs into logical pages (~500 words each)."""
    try:
        from docx import Document

        doc   = Document(io.BytesIO(data))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return _chunk_into_pages(paras, words_per_page=400)
    except Exception as e:
        return [f"[Error extracting DOCX: {str(e)}]"]


def extract_pptx(data: bytes) -> List[str]:
    """Extract text from PPTX bytes. One slide = one page."""
    try:
        from pptx import Presentation

        prs   = Presentation(io.BytesIO(data))
        pages = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                pages.append("\n\n".join(texts))
        return pages if pages else ["[No text content found in this presentation.]"]
    except Exception as e:
        return [f"[Error extracting PPTX: {str(e)}]"]


def extract_text(data: bytes, file_type: str) -> List[str]:
    """
    Public interface. Dispatches to the correct extractor.
    file_type: 'pdf' | 'docx' | 'pptx'
    Returns list of page strings (one string per page).
    """
    file_type = file_type.lower().strip(".")
    if file_type == "pdf":
        return extract_pdf(data)
    elif file_type in ("docx", "doc"):
        return extract_docx(data)
    elif file_type in ("pptx", "ppt"):
        return extract_pptx(data)
    else:
        return [f"[Unsupported file type: {file_type}]"]


def _chunk_into_pages(paragraphs: List[str], words_per_page: int = 400) -> List[str]:
    """Group paragraphs into pages of approximately words_per_page words."""
    pages        = []
    current_page = []
    current_words = 0

    for para in paragraphs:
        word_count = len(para.split())
        if current_words + word_count > words_per_page and current_page:
            pages.append("\n\n".join(current_page))
            current_page  = [para]
            current_words = word_count
        else:
            current_page.append(para)
            current_words += word_count

    if current_page:
        pages.append("\n\n".join(current_page))

    return pages if pages else ["[Document appears to be empty.]"]
