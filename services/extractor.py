"""
Document text extractor
-----------------------
Supports PDF (pdfplumber), DOCX (python-docx).
Returns a list of strings — one per "page".
For DOCX, we group every ~500 words into a synthetic page.
"""
from __future__ import annotations
import io


def extract_pages(file_bytes: bytes, file_type: str) -> list[str]:
    """Return a list of page text strings."""
    ft = file_type.lower()
    if ft == "pdf":
        return _extract_pdf(file_bytes)
    if ft == "docx":
        return _extract_docx(file_bytes)
    if ft == "pptx":
        return _extract_pptx(file_bytes)
    raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf(data: bytes) -> list[str]:
    import pdfplumber  # type: ignore
    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = (page.extract_text() or "").strip()
            pages.append(text if text else "[No text content on this page]")
    return pages or ["[Empty document]"]


def _extract_docx(data: bytes) -> list[str]:
    from docx import Document  # type: ignore
    doc    = Document(io.BytesIO(data))
    pages: list[str] = []
    current: list[str] = []
    word_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        current.append(text)
        word_count += len(text.split())
        if word_count >= 500:
            pages.append("\n\n".join(current))
            current = []
            word_count = 0

    if current:
        pages.append("\n\n".join(current))

    return pages or ["[Empty document]"]


def _extract_pptx(data: bytes) -> list[str]:
    """Extract text from each slide of a PowerPoint file."""
    try:
        from pptx import Presentation  # type: ignore
        prs    = Presentation(io.BytesIO(data))
        pages: list[str] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            pages.append("\n".join(texts) if texts else "[Slide with no text]")
        return pages or ["[Empty presentation]"]
    except ImportError:
        return ["[PPTX extraction requires python-pptx — add to requirements.txt]"]
